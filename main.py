import tkinter as tk
from tkinter import filedialog
import os
from pptx import Presentation
from PIL import Image
from pptx.util import Inches,Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


root = tk.Tk()
root.withdraw()
file_paths = filedialog.askopenfilenames(title="Select photos", filetypes=[("Image files", "*.png *.jpg *.jpeg *.JPG *.PNG *.JPEG"),("Video files", "*.mp4 *.mov *.MOV *.MP4")])  #


# Create a new PowerPoint presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])

# Insert each photo or video on a new slide
for file_path in file_paths:
    # Get the file extension
    file_name = os.path.splitext(file_path)[0].split("/")[-1]
    print(file_name)
    file_ext = os.path.splitext(file_path)[1].lower()
    rotate = False

    if file_ext in ('.png', '.jpg', '.jpeg'):
        # Insert the photo on a new slide
        image = Image.open(file_path)
        exif = image.getexif()
        try:
            rotate = exif[274]
        except:
            pass

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pic = slide.shapes.add_picture(file_path,left=0,top=0)

        slide_width = prs.slide_width
        slide_height = prs.slide_height
        pic_width = pic.width
        pic_height = pic.height
        slide_ratio = slide_width / slide_height
        pic_ratio = pic_width / pic_height

        if slide_ratio > pic_ratio:
            # Slide is wider than photo, so scale the photo to fill the width of the slide
            width_scale = slide_width / pic_width
            pic.width = slide_width
            pic.height = int(pic_height * width_scale)

        else:
            # Slide is taller than photo, so scale the photo to fill the height of the slide
            height_scale = slide_height / pic_height
            pic.width = int(pic_width * height_scale)
            pic.height = slide_height

        # Center the picture in the middle of the slide
        left = int((slide_width - pic.width) / 2)
        top = int((slide_height - pic.height) / 2)
        pic.left = left
        pic.top = top


    # VIDEO IS NOT WORKING PROPERLY
    elif file_ext in ('.mp4', '.mov' ):
    #     from videoprops import get_video_properties
    #     import cv2
    #     # Import everything needed to edit video clips
    #     props = get_video_properties(file_path)
    #     vidcap = cv2.VideoCapture(file_path)
    #     # get total number of frames
    #     # set frame position
    #     vidcap.set(cv2.CAP_PROP_POS_FRAMES, 10)
    #     success, image = vidcap.read()
    #     print(type(image))
    #     print(success)
    #     # Insert the video on a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
        text_frame = textbox.text_frame
        text_frame.text = file_name
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Inches(1)





    #     pic = slide.shapes.add_movie(file_path, top=0,left=0, width=Pt(0.75*props['width']), height=Pt(0.75*props['height']),poster_frame_image=image)
    #
    #     if props['width'] < props['height']:
    #         pic.rotation = -90
        # Calculate the scaling factor needed to fill the slide with the photo

# Save the PowerPoint presentation
prs.save('output.pptx')
